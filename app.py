# app.py – Streamlit PPTX generator for IR derivatives deck
# -*- coding: utf-8 -*-
import re
import datetime as dt
from pathlib import Path
import os

import requests
import pandas as pd
import matplotlib.pyplot as plt
import streamlit as st
from pptx import Presentation

# =========================
# CONFIG
# =========================
# Template will be stored in the repo root alongside app.py
TEMPLATE_PPTX_PATH = Path("Derivatives_Presentation_Template.pptx")

# Output directory (relative, works on Streamlit Cloud)
OUTPUT_DIR = Path("output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

SOFR_GRAPH_PLACEHOLDER_KEY  = "SOFR_GRAPH"
DGS10_GRAPH_PLACEHOLDER_KEY = "DGS10_GRAPH"

SOFR_CHART_PATH  = OUTPUT_DIR / "sofr_chart.png"
DGS10_CHART_PATH = OUTPUT_DIR / "dgs10_chart.png"

FRED_SOFR_SERIES  = "SOFR"
FRED_DGS10_SERIES = "DGS10"
FRED_LOOKBACK_DAYS_DEFAULT = 365 * 2

# Default data (you’ll override via UI)
DATA_DEFAULT = {
    "Presentation Date": "December 17, 2025",
    "Marketer Name": "Gil De La Cruz Vazquez",
    "Marketer Title": "Sales Execution Analyst",
    "Business Division": "Capital Markets Division",
    "Client Legal Name": "Acme Manufacturing, Inc.",

    "Loan Amount": "$5,000,000",
    "Effective Date": "January 5, 2026",
    "Swap Term": "7 Years",
    "Amortization": "25 Years",
    "Adjustable Rate Index": "1-Month Term SOFR",
    "Adjustable Rate Spread": "2.108%",
    "Payment Frequency": "Monthly",
    "Interest Accrual Basis": "Actual/360",

    "Fixed Rate Payer": "Acme Manufacturing, Inc.",
    "Float Rate Payer": "JPMorgan Chase Bank, N.A.",
    "Swap Rate": "4.052%",

    "Index Rate": "0.30%",
    "index rate": "0.30%",
    "Client Position PAYs / RECs": "PAYs",
    "Client Position PAYs/RECs": "PAYs",
    "Client Leg fixed/float": "fixed",
    "Bank Position PAYs / RECs": "RECs",
    "Bank Position PAYs/RECs": "RECs",
}

PLACEHOLDER_RE = re.compile(r"\{\{([^{}]+)\}\}")


# =========================
# Core helpers
# =========================
def get_fred_api_key() -> str:
    """Prefer st.secrets in prod; fall back to env var for local dev."""
    try:
        key = st.secrets.get("FRED_API_KEY", "")
    except Exception:
        key = ""
    if not key:
        key = os.getenv("FRED_API_KEY", "")
    return key.strip()


def safe_filename(s: str, max_len: int = 80) -> str:
    s = (s or "").strip()
    s = re.sub(r'[\\/:*?"<>|]+', "", s)
    s = re.sub(r"\s+", " ", s)
    s = s.strip(" .")
    return (s[:max_len].strip() or "Client")


def build_output_path(data: dict) -> Path:
    client = safe_filename(data.get("Client Legal Name", "Client"))
    today = dt.date.today().strftime("%Y-%m-%d")
    return OUTPUT_DIR / f"{client} - {today}.pptx"


def pct_to_float(pct: str) -> float:
    if pct is None:
        raise ValueError("Percent value is None")
    s = str(pct).strip()
    if s.endswith("%"):
        return float(s[:-1].strip()) / 100.0
    return float(s) / 100.0


def compute_derived_fields(data: dict) -> None:
    _index = pct_to_float(data.get("Index Rate", "0%"))
    _spread = pct_to_float(data.get("Adjustable Rate Spread", "0%"))
    _loan_adj = _index + _spread
    data["Loan Adj Rate"] = f"{_loan_adj * 100:.3f}%"


def fred_get_series_observations(series_id: str, api_key: str, observation_start: str) -> pd.Series:
    if not api_key:
        raise ValueError("Missing FRED_API_KEY.")
    url = "https://api.stlouisfed.org/fred/series/observations"
    params = {
        "series_id": series_id,
        "api_key": api_key,
        "file_type": "json",
        "observation_start": observation_start,
    }
    r = requests.get(url, params=params, timeout=30)
    r.raise_for_status()
    payload = r.json()
    obs = payload.get("observations", [])
    if not obs:
        raise RuntimeError(f"No observations for {series_id}")
    rows = []
    for o in obs:
        date_str = o["date"]
        val_str = o["value"]
        if val_str == ".":
            continue
        rows.append((pd.to_datetime(date_str), float(val_str)))
    s = pd.Series(dict(rows)).sort_index()
    s.name = series_id
    return s


def build_line_chart_from_fred(
    series_id: str,
    api_key: str,
    title: str,
    ylabel: str,
    out_png: Path,
    lookback_days: int | None = 365,
    observation_start: str | None = None,
) -> dict:
    if observation_start:
        start_date = observation_start
    else:
        if lookback_days is None:
            raise ValueError("Either observation_start or lookback_days must be provided")
        start_date = (dt.date.today() - dt.timedelta(days=lookback_days)).strftime("%Y-%m-%d")

    s = fred_get_series_observations(series_id, api_key, start_date)

    plt.figure(figsize=(10, 5))
    plt.plot(s.index, s.values, label=series_id)
    plt.title(title)
    plt.ylabel(ylabel)
    plt.xlabel("Date")
    plt.grid(True, alpha=0.3)
    plt.tight_layout()

    out_png.parent.mkdir(parents=True, exist_ok=True)
    plt.savefig(out_png, dpi=200)
    plt.close()

    latest = {}
    if not s.empty:
        latest = {
            f"{series_id} Latest": f"{s.iloc[-1]:.3f}%",
            f"{series_id} As Of": s.index.max().strftime("%Y-%m-%d"),
        }
    return latest


def replace_placeholders(text: str, data: dict) -> str:
    if not text:
        return text

    def repl(m: re.Match) -> str:
        key = m.group(1).strip()
        if key in data and data[key] is not None:
            return str(data[key])
        return m.group(0)

    return PLACEHOLDER_RE.sub(repl, text)


def iter_all_shapes(slide):
    for shape in slide.shapes:
        yield shape
        if getattr(shape, "shape_type", None) == 6:
            for subshape in shape.shapes:
                yield subshape


def paragraph_full_text(paragraph) -> str:
    return "".join(run.text for run in paragraph.runs)


def set_paragraph_text_preserve_runs(paragraph, new_text: str) -> None:
    if not paragraph.runs:
        paragraph.text = new_text
        return
    paragraph.runs[0].text = new_text
    for r in paragraph.runs[1:]:
        r.text = ""


def replace_in_text_frame(text_frame, data: dict) -> None:
    for paragraph in text_frame.paragraphs:
        full = paragraph_full_text(paragraph)
        replaced = replace_placeholders(full, data)
        if replaced != full:
            set_paragraph_text_preserve_runs(paragraph, replaced)


def replace_in_table(table, data: dict) -> None:
    for row in table.rows:
        for cell in row.cells:
            tf = cell.text_frame
            if tf is None:
                continue
            replace_in_text_frame(tf, data)


def collect_used_keys(prs: Presentation) -> set[str]:
    used = set()
    for slide in prs.slides:
        for shape in iter_all_shapes(slide):
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    used |= {
                        m.group(1).strip()
                        for m in PLACEHOLDER_RE.finditer(paragraph_full_text(p))
                    }
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            used |= {
                                m.group(1).strip()
                                for m in PLACEHOLDER_RE.finditer(paragraph_full_text(p))
                            }
    return used


def collect_unresolved(prs: Presentation) -> dict:
    unresolved = {}
    for si, slide in enumerate(prs.slides, start=1):
        keys = set()
        for shape in iter_all_shapes(slide):
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    keys |= {
                        m.group(1).strip()
                        for m in PLACEHOLDER_RE.finditer(paragraph_full_text(p))
                    }
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            keys |= {
                                m.group(1).strip()
                                for m in PLACEHOLDER_RE.finditer(paragraph_full_text(p))
                            }
        if keys:
            unresolved[si] = sorted(keys)
    return unresolved


def insert_image_by_placeholder_on_slide(
    prs: Presentation,
    slide_index_1based: int,
    placeholder_key: str,
    image_path: Path
) -> int:
    token = f"{{{{{placeholder_key}}}}}"
    if slide_index_1based < 1 or slide_index_1based > len(prs.slides):
        return 0
    slide = prs.slides[slide_index_1based - 1]
    inserted = 0
    shapes = list(slide.shapes)
    for shape in shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            full_text = shape.text_frame.text or ""
            if token in full_text:
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                sp = shape._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
                inserted += 1
    return inserted


def generate_pptx(data: dict, fred_api_key: str, fred_lookback_days: int) -> Path:
    """Run the pipeline and return path to generated PPTX."""
    compute_derived_fields(data)

    latest_sofr = build_line_chart_from_fred(
        series_id=FRED_SOFR_SERIES,
        api_key=fred_api_key,
        title="Secured Overnight Financing Rate (SOFR)",
        ylabel="Percent",
        out_png=SOFR_CHART_PATH,
        lookback_days=fred_lookback_days,
        observation_start=None,
    )
    latest_dgs10 = build_line_chart_from_fred(
        series_id=FRED_DGS10_SERIES,
        api_key=fred_api_key,
        title="10-Year Treasury Yield (DGS10)",
        ylabel="Percent",
        out_png=DGS10_CHART_PATH,
        observation_start="1984-01-01",
        lookback_days=None,
    )

    data.update(latest_sofr)
    data.update(latest_dgs10)

    if not TEMPLATE_PPTX_PATH.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PPTX_PATH}")

    prs = Presentation(str(TEMPLATE_PPTX_PATH))
    out_path = build_output_path(data)

    used_keys = collect_used_keys(prs)
    missing = sorted(k for k in used_keys if k not in data)
    if missing:
        st.warning("Missing values for placeholders: " + ", ".join(missing))

    for slide in prs.slides:
        for shape in iter_all_shapes(slide):
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                replace_in_text_frame(shape.text_frame, data)
            if getattr(shape, "has_table", False) and shape.has_table:
                replace_in_table(shape.table, data)

    insert_image_by_placeholder_on_slide(
        prs, slide_index_1based=2,
        placeholder_key=SOFR_GRAPH_PLACEHOLDER_KEY,
        image_path=SOFR_CHART_PATH
    )
    insert_image_by_placeholder_on_slide(
        prs, slide_index_1based=3,
        placeholder_key=DGS10_GRAPH_PLACEHOLDER_KEY,
        image_path=DGS10_CHART_PATH
    )

    prs.save(str(out_path))
    return out_path


# =========================
# STREAMLIT UI
# =========================
def main():
    st.title("Interest Rate Derivatives Deck Generator")

    st.caption("Template: `Derivatives_Presentation_Template.pptx` (bundled in repo)")

    fred_key = get_fred_api_key()
    if not fred_key:
        st.error("FRED_API_KEY not set in Streamlit secrets or environment.")
        st.stop()

    fred_lookback_days = st.sidebar.number_input(
        "SOFR lookback days",
        min_value=30,
        max_value=3650,
        value=FRED_LOOKBACK_DAYS_DEFAULT,
        step=30,
    )

    st.header("Client & Deal Inputs")

    with st.form("inputs"):
        col1, col2 = st.columns(2)

        with col1:
            pres_date      = st.text_input("Presentation Date", DATA_DEFAULT["Presentation Date"])
            marketer_name  = st.text_input("Marketer Name", DATA_DEFAULT["Marketer Name"])
            marketer_title = st.text_input("Marketer Title", DATA_DEFAULT["Marketer Title"])
            biz_div        = st.text_input("Business Division", DATA_DEFAULT["Business Division"])
            client_legal   = st.text_input("Client Legal Name", DATA_DEFAULT["Client Legal Name"])
            loan_amount    = st.text_input("Loan Amount", DATA_DEFAULT["Loan Amount"])
            effective_date = st.text_input("Effective Date", DATA_DEFAULT["Effective Date"])
            swap_term      = st.text_input("Swap Term", DATA_DEFAULT["Swap Term"])
            amortization   = st.text_input("Amortization", DATA_DEFAULT["Amortization"])

        with col2:
            adj_index   = st.text_input("Adjustable Rate Index", DATA_DEFAULT["Adjustable Rate Index"])
            adj_spread  = st.text_input("Adjustable Rate Spread", DATA_DEFAULT["Adjustable Rate Spread"])
            pay_freq    = st.text_input("Payment Frequency", DATA_DEFAULT["Payment Frequency"])
            accrual     = st.text_input("Interest Accrual Basis", DATA_DEFAULT["Interest Accrual Basis"])
            fixed_payer = st.text_input("Fixed Rate Payer", DATA_DEFAULT["Fixed Rate Payer"])
            float_payer = st.text_input("Float Rate Payer", DATA_DEFAULT["Float Rate Payer"])
            swap_rate   = st.text_input("Swap Rate", DATA_DEFAULT["Swap Rate"])
            index_rate  = st.text_input("Index Rate", DATA_DEFAULT["Index Rate"])
            client_pos  = st.text_input("Client Position PAYs/RECs", DATA_DEFAULT["Client Position PAYs/RECs"])
            client_leg  = st.text_input("Client Leg fixed/float", DATA_DEFAULT["Client Leg fixed/float"])
            bank_pos    = st.text_input("Bank Position PAYs/RECs", DATA_DEFAULT["Bank Position PAYs/RECs"])

        submitted = st.form_submit_button("Generate Presentation")

    if submitted:
        data = DATA_DEFAULT.copy()
        data.update({
            "Presentation Date": pres_date,
            "Marketer Name": marketer_name,
            "Marketer Title": marketer_title,
            "Business Division": biz_div,
            "Client Legal Name": client_legal,
            "Loan Amount": loan_amount,
            "Effective Date": effective_date,
            "Swap Term": swap_term,
            "Amortization": amortization,
            "Adjustable Rate Index": adj_index,
            "Adjustable Rate Spread": adj_spread,
            "Payment Frequency": pay_freq,
            "Interest Accrual Basis": accrual,
            "Fixed Rate Payer": fixed_payer,
            "Float Rate Payer": float_payer,
            "Swap Rate": swap_rate,
            "Index Rate": index_rate,
            "index rate": index_rate,
            "Client Position PAYs / RECs": client_pos,
            "Client Position PAYs/RECs": client_pos,
            "Client Leg fixed/float": client_leg,
            "Bank Position PAYs / RECs": bank_pos,
            "Bank Position PAYs/RECs": bank_pos,
        })

        try:
            out_path = generate_pptx(
                data=data,
                fred_api_key=fred_key,
                fred_lookback_days=int(fred_lookback_days),
            )
        except Exception as e:
            st.error(f"Error generating PPTX: {e}")
            return

        st.success("Presentation generated.")
        with open(out_path, "rb") as f:
            st.download_button(
                label="Download PPTX",
                data=f.read(),
                file_name=out_path.name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )


if __name__ == "__main__":
    main()
